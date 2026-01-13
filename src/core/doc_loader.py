import fitz  # PyMuPDF
from pptx import Presentation
import os

def load_documents(file_paths):
    chunks = []
    
    for path in file_paths:
        filename = os.path.basename(path)

        # PDF Extraction
        if path.endswith(".pdf"):
            doc = fitz.open(path)
            for i, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    chunks.append({
                        "content": text,
                        "metadata": {
                            "filename": filename,
                            "page_no": i+1,
                            "type": "text"
                        }
                    })

        # PPT Extraction
        elif path.endswith(".pptx"):
            prs = Presentation(path)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            slide_text.append(para.text.strip())
                combined = "\n".join(slide_text).strip()
                if combined:
                    chunks.append({
                        "content": combined,
                        "metadata": {
                            "filename": filename,
                            "slide_no": i+1,
                            "type": "text"
                        }
                    })

        # Ignore unsupported formats for now
        else:
            print(f"Skipping unsupported file: {filename}")

    return chunks
